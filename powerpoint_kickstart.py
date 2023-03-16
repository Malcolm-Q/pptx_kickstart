# import and mount
import requests
from pptx import Presentation
from pptx.util import Inches
import openai
from os import environ

# This is a script that save a powerpoint file given a prompt on what you want the presentation to be about
# the powerpoint will have basic slides showing the structure of the presentation you can easily tweak further
def powerpoint(idea = '', verbose = True, key='openai_key', path =''):
    '''Uses chatgpt to create a basic powerpoint with a structure and ideas on a given topic.
        params:
            key (str) : name of environment variable for openai API key. Default = 'openai_key',
            idea (str) : the prompt to build a presentation about,
            verbose (bool) : whether to print updates on the process. default = true,
            path (str) : path to save in (can be left blank),

        returns:
            nothing, it saves a file in working directory or path provided.'''

    # get content from ChatGPT turbo
    if idea == '': idea = input('Enter a topic:')
    if verbose: print(f'Loading environment variable "{key}"')
    key = environ[key]
    openai.api_key=key

    if verbose: print('\nGenerating presentation...')
    request = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that builds outlines and notes for powerpoint presentations. Your responses can start with 'Slide ', or '- ' for details on title slides"},
            {"role": "user", "content": f"Please build some slides for a presentation on {idea}"}
        ]
    )
    if verbose: print(f'Presentation data has been generated!\nCleaning and building Powerpoint.')

    # we do a little cleaning
    words = request['choices'][0]['message']['content']
    sentences = words.split('\n')
    clean = []
    for sentence in sentences:
        # cut empty values that will break indexes below
        if(sentence == ''): continue
        # make sure it's listening to system prompt
        if sentence[0] in ['S','-']:
            if(sentence[0:4] == 'Sure'): continue
            clean.append(sentence)

    # initialize Presentation and choose layout
    X = Presentation()
    Layout = X.slide_layouts[1]
    slide = None
    ypos = 6

    if verbose: print(f'Presentation has {len(clean)} bodies.')

    for sentence in clean:
        # we could likely fine tune gpt or use a more strict system prompt to not have to do
        # this wonky handling of its message but oh well this works.
        
        if sentence[0] == '-':
            if slide.placeholders[1].text == '': slide.placeholders[1].text = sentence[2:]
            else:
            # if subtitle / note field is already assigned create and fill a new one.
                new_textbox = slide.shapes.add_textbox(Inches(1.5), Inches(ypos), Inches(8), Inches(5))
                new_textbox.text = sentence[2:]
                ypos-=1
      
        elif sentence[0] == 'S':
            ypos = 6
            slide = X.slides.add_slide(Layout)
            slide.shapes.title.text = sentence.split(':')[1]

    # saving
    if path == '':
        X.save(f"{idea}.pptx")
        if verbose: print(f'\nDone!\nPowerpoint has been saved to working directory as "{idea}.pptx"')
    else:
        X.save(f"{path}/{idea}.pptx")
        if verbose: print(f'\nDone!\nPowerpoint has been saved path: {path}/{idea}.pptx')
    

if __name__ == '__main__':
  print('~MODULE EXECUTED~')
  powerpoint()